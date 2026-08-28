"""Light text extractors for zipped office formats (docx, xlsx, pptx).

Why not SimpleDirectoryReader: its per-format readers live in the optional
``llama-index-readers-file`` distribution, which this deployment does not
install — and its documented fallback for an unknown extension is to read the
raw bytes as UTF-8 text. For a ``.docx`` (a zip) that produced ``PK\\x03…``
garbage which the binary-content guard then rejected, so every Word upload
failed ingestion while the UI offered the type. These extractors are small,
depend only on libraries the extra already carries (docx2txt, python-pptx,
openpyxl), and give every unit a real ``page_label`` (sheet name, slide
number) a citation can point at.

Only formats the generic reader actually garbles get a handler: plain-text
formats (``.csv``, ``.txt``, ``.md``) ingest correctly through the fallback
text read and deliberately have none.

Modularity contract (see ``docs/architecture/visual-ingestion.md``): one
handler per format, keyed by extension, returning plain ``Document`` objects —
adding a format is one function plus one dict entry, and the adapter knows
nothing about any of them.
"""

from __future__ import annotations

import logging
from collections.abc import Callable
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

# Bounds: a schedule/BoQ spreadsheet can carry tens of thousands of rows; the
# text form exists for retrieval, not for round-tripping, so cap per unit and
# say so in the emitted text rather than silently truncating.
MAX_TABLE_ROWS = 1000
MAX_TABLE_COLS = 60
MAX_CELL_CHARS = 500


def _markdown_table(rows: list[list[str]], truncated: bool) -> str:
    """Render rows as the same markdown shape the PDF table extractor emits."""
    if not rows:
        return ""
    header, *body = rows
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(["---"] * len(header)) + " |",
    ]
    for row in body:
        padded = row + [""] * (len(header) - len(row))
        lines.append("| " + " | ".join(padded[: len(header)]) + " |")
    if truncated:
        lines.append("")
        lines.append(f"[Tabelle gekürzt: nur die ersten {MAX_TABLE_ROWS} Zeilen indexiert]")
    return "\n".join(lines)


def _clean_cell(value: Any) -> str:
    if value is None:
        return ""
    return str(value).replace("|", "\\|").replace("\n", " ").strip()[:MAX_CELL_CHARS]


def _extract_docx(file_path: str) -> list[tuple[str, str]]:
    """Word text via docx2txt: one unit, no page mapping (docx has no pages)."""
    import docx2txt

    text = (docx2txt.process(file_path) or "").strip()
    return [("1", text)] if text else []


def _extract_xlsx(file_path: str) -> list[tuple[str, str]]:
    """One unit per worksheet, rendered as a markdown table, labeled by sheet name."""
    from openpyxl import load_workbook

    units: list[tuple[str, str]] = []
    workbook = load_workbook(file_path, read_only=True, data_only=True)
    try:
        for sheet in workbook.worksheets:
            rows: list[list[str]] = []
            truncated = False
            for row_index, row in enumerate(sheet.iter_rows(values_only=True)):
                if row_index >= MAX_TABLE_ROWS:
                    truncated = True
                    break
                cells = [_clean_cell(cell) for cell in row[:MAX_TABLE_COLS]]
                if any(cells):
                    rows.append(cells)
            if rows:
                units.append((sheet.title, f"Tabellenblatt „{sheet.title}“\n\n{_markdown_table(rows, truncated)}"))
    finally:
        workbook.close()
    return units


def _extract_pptx(file_path: str) -> list[tuple[str, str]]:
    """One unit per slide: text frames, tables and speaker notes."""
    from pptx import Presentation

    units: list[tuple[str, str]] = []
    presentation = Presentation(file_path)
    for number, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
            if getattr(shape, "has_table", False):
                rows = [[_clean_cell(cell.text) for cell in row.cells] for row in shape.table.rows]
                rows = [row for row in rows if any(row)]
                if rows:
                    parts.append(_markdown_table(rows[:MAX_TABLE_ROWS], len(rows) > MAX_TABLE_ROWS))
        notes = getattr(slide, "notes_slide", None)
        if notes is not None and notes.notes_text_frame is not None:
            note_text = notes.notes_text_frame.text.strip()
            if note_text:
                parts.append(f"Notizen: {note_text}")
        if parts:
            units.append((str(number), f"Folie {number}\n\n" + "\n\n".join(parts)))
    return units


#: Extension → handler. Adding a format is one function plus one entry here;
#: macro-enabled variants share their sibling's handler (same XML inside).
_HANDLERS: dict[str, Callable[[str], list[tuple[str, str]]]] = {
    ".docx": _extract_docx,
    ".docm": _extract_docx,
    ".xlsx": _extract_xlsx,
    ".xlsm": _extract_xlsx,
    ".pptx": _extract_pptx,
    ".pptm": _extract_pptx,
}

#: Formats this module handles (advertisable to capability probes/tests).
SUPPORTED_EXTENSIONS = frozenset(_HANDLERS)


def extract_office_documents(file_path: str, file_name: str, file_size: int) -> list[Any] | None:
    """Extract a known office format into LlamaIndex ``Document`` objects.

    Returns ``None`` when the extension (of the ORIGINAL ``file_name`` first,
    the temp ``file_path`` second) is not one this module handles — the caller
    then falls through to its generic reader. A handled-but-empty file returns
    ``[]`` so the pipeline fails that file with its normal "no content" path
    rather than indexing nothing silently. Extraction errors propagate: the
    per-file error handling in ``_run_ingestion`` reports them per file.
    """
    extension = (Path(file_name).suffix or Path(file_path).suffix).lower()
    handler = _HANDLERS.get(extension)
    if handler is None:
        return None

    from llama_index.core import Document

    units = handler(file_path)
    logger.info("Office extraction (%s): %d unit(s) from %s", extension, len(units), file_name)
    return [
        Document(
            text=text,
            metadata={
                "file_name": file_name,
                "file_size": file_size,
                "page_label": label,
                "content_type": "table" if extension in {".xlsx", ".xlsm"} else "text",
            },
        )
        for label, text in units
    ]

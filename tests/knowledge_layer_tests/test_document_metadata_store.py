"""Unit tests for document summarization storage.

Tests cover:
- DocumentMetadataStore SQLAlchemy-based storage (SQLite)
- Factory functions (register_summary, get_available_documents, etc.)
- AvailableDocument model
- URL normalization for database connections
"""

import tempfile
from pathlib import Path

import pytest

from aiq_agent.knowledge.document_metadata_store import DocumentMetadataStore
from aiq_agent.knowledge.document_metadata_store import _normalize_db_url
from aiq_agent.knowledge.schema import AvailableDocument

# =============================================================================
# URL Normalization Tests
# =============================================================================


class TestNormalizeDbUrl:
    """Tests for the _normalize_db_url helper function."""

    def test_sqlite_url_async_mode(self):
        """Test SQLite URL normalization for async mode."""
        url = "sqlite:///./summaries.db"
        result = _normalize_db_url(url, async_mode=True)
        assert result == "sqlite+aiosqlite:///./summaries.db"

    def test_sqlite_url_sync_mode(self):
        """Test SQLite URL normalization for sync mode."""
        url = "sqlite:///./summaries.db"
        result = _normalize_db_url(url, async_mode=False)
        assert result == "sqlite:///./summaries.db"

    def test_sqlite_already_async(self):
        """Test SQLite URL that already has async driver."""
        url = "sqlite+aiosqlite:///./summaries.db"
        result = _normalize_db_url(url, async_mode=True)
        assert result == "sqlite+aiosqlite:///./summaries.db"

    def test_postgresql_url_async_mode(self):
        """Test PostgreSQL URL normalization for async mode."""
        url = "postgresql://user:pass@localhost:5432/db"  # pragma: allowlist secret
        result = _normalize_db_url(url, async_mode=True)
        assert result == "postgresql+psycopg://user:pass@localhost:5432/db"  # pragma: allowlist secret

    def test_postgresql_url_sync_mode(self):
        """Test PostgreSQL URL normalization for sync mode."""
        url = "postgresql://user:pass@localhost:5432/db"  # pragma: allowlist secret
        result = _normalize_db_url(url, async_mode=False)
        assert result == "postgresql+psycopg://user:pass@localhost:5432/db"  # pragma: allowlist secret

    def test_postgres_shorthand_url(self):
        """Test postgres:// shorthand URL normalization."""
        url = "postgres://user:pass@localhost:5432/db"  # pragma: allowlist secret
        result = _normalize_db_url(url, async_mode=True)
        assert result == "postgresql+psycopg://user:pass@localhost:5432/db"  # pragma: allowlist secret

    def test_postgresql_with_existing_driver(self):
        """Test PostgreSQL URL with existing driver gets normalized."""
        url = "postgresql+asyncpg://user:pass@localhost:5432/db"  # pragma: allowlist secret
        result = _normalize_db_url(url, async_mode=True)
        assert "psycopg" in result

    def test_asyncpg_ssl_param_translated_to_sslmode(self):
        """asyncpg's `ssl` query param must become psycopg's `sslmode` on driver rewrite."""
        url = "postgresql+asyncpg://user:pass@localhost:5432/db?ssl=require"  # pragma: allowlist secret
        result = _normalize_db_url(url, async_mode=False)
        assert result == "postgresql+psycopg://user:pass@localhost:5432/db?sslmode=require"  # pragma: allowlist secret

    def test_asyncpg_ssl_param_translated_mid_query(self):
        """The translation also applies when `ssl` is not the first query parameter."""
        url = "postgresql+asyncpg://u:p@localhost:5432/db?application_name=x&ssl=require"  # pragma: allowlist secret
        result = _normalize_db_url(url, async_mode=False)
        expected = (
            "postgresql+psycopg://u:p@localhost:5432/db"  # pragma: allowlist secret
            "?application_name=x&sslmode=require"
        )
        assert result == expected

    def test_asyncpg_ssl_boolean_values_mapped(self):
        """asyncpg's boolean ssl values map to libpq's require/disable."""
        url = "postgresql+asyncpg://u:p@localhost:5432/db?ssl=true"  # pragma: allowlist secret
        result = _normalize_db_url(url)
        assert result == "postgresql+psycopg://u:p@localhost:5432/db?sslmode=require"  # pragma: allowlist secret
        url = "postgresql+asyncpg://u:p@localhost:5432/db?ssl=false"  # pragma: allowlist secret
        result = _normalize_db_url(url)
        assert result == "postgresql+psycopg://u:p@localhost:5432/db?sslmode=disable"  # pragma: allowlist secret

    def test_sslmode_param_passes_through(self):
        """A libpq-style `sslmode` param is already correct and stays untouched."""
        url = "postgresql+psycopg://user:pass@localhost:5432/db?sslmode=require"  # pragma: allowlist secret
        result = _normalize_db_url(url, async_mode=True)
        assert result == "postgresql+psycopg://user:pass@localhost:5432/db?sslmode=require"  # pragma: allowlist secret

    def test_unknown_url_passthrough(self):
        """Test unknown database URLs pass through unchanged."""
        url = "mysql://user:pass@localhost/db"  # pragma: allowlist secret
        result = _normalize_db_url(url, async_mode=True)
        assert result == url


# =============================================================================
# AvailableDocument Model Tests
# =============================================================================


class TestAvailableDocument:
    """Tests for the AvailableDocument Pydantic model."""

    def test_create_with_summary(self):
        """Test creating AvailableDocument with a summary."""
        doc = AvailableDocument(file_name="test.pdf", summary="A test document.")
        assert doc.file_name == "test.pdf"
        assert doc.summary == "A test document."

    def test_create_without_summary(self):
        """Test creating AvailableDocument without a summary."""
        doc = AvailableDocument(file_name="test.pdf")
        assert doc.file_name == "test.pdf"
        assert doc.summary is None

    def test_model_dump(self):
        """Test model serialization to dict."""
        doc = AvailableDocument(file_name="report.pdf", summary="Financial report.")
        data = doc.model_dump()
        assert data == {
            "file_name": "report.pdf",
            "summary": "Financial report.",
            "tags": None,
            "doc_class": None,
            "display_title": None,
        }

    def test_model_dump_without_summary(self):
        """Test model serialization without summary."""
        doc = AvailableDocument(file_name="report.pdf")
        data = doc.model_dump()
        assert data == {
            "file_name": "report.pdf",
            "summary": None,
            "tags": None,
            "doc_class": None,
            "display_title": None,
        }

    def test_model_validate(self):
        """Test model creation from dict."""
        data = {"file_name": "doc.pdf", "summary": "Test summary"}
        doc = AvailableDocument.model_validate(data)
        assert doc.file_name == "doc.pdf"
        assert doc.summary == "Test summary"
        assert doc.tags is None

    def test_create_with_tags(self):
        """Test creating AvailableDocument with controlled tags."""
        doc = AvailableDocument(file_name="plan.pdf", summary="A floor plan.", tags=["Grundriss", "Brandschutz"])
        assert doc.tags == ["Grundriss", "Brandschutz"]
        # Tags survive a serialize/round-trip through model_dump.
        assert AvailableDocument.model_validate(doc.model_dump()).tags == ["Grundriss", "Brandschutz"]


# =============================================================================
# DocumentMetadataStore Tests
# =============================================================================


class TestDocumentMetadataStore:
    """Tests for the DocumentMetadataStore SQLAlchemy-based storage."""

    @pytest.fixture
    def temp_db(self):
        """Create a temporary SQLite database for testing."""
        with tempfile.TemporaryDirectory() as tmpdir:
            db_path = Path(tmpdir) / "test_summaries.db"
            db_url = f"sqlite:///{db_path}"
            yield db_url
            DocumentMetadataStore.dispose_engine(db_url)

    @pytest.fixture
    def store(self, temp_db):
        """Create a DocumentMetadataStore instance with temp database."""
        return DocumentMetadataStore(temp_db)

    def test_store_initialization(self, temp_db):
        """Test DocumentMetadataStore initializes correctly."""
        store = DocumentMetadataStore(temp_db)
        assert store.db_url == temp_db

    def test_register_summary(self, store):
        """Test registering a document summary."""
        store.register("test_collection", "doc1.pdf", "This is a test summary.")
        docs = store.get_all("test_collection")
        assert len(docs) == 1
        assert docs[0].file_name == "doc1.pdf"
        assert docs[0].summary == "This is a test summary."

    def test_register_multiple_summaries(self, store):
        """Test registering multiple document summaries."""
        store.register("collection1", "doc1.pdf", "Summary 1")
        store.register("collection1", "doc2.pdf", "Summary 2")
        store.register("collection1", "doc3.pdf", "Summary 3")

        docs = store.get_all("collection1")
        assert len(docs) == 3
        filenames = {doc.file_name for doc in docs}
        assert filenames == {"doc1.pdf", "doc2.pdf", "doc3.pdf"}

    def test_register_updates_existing(self, store):
        """Test registering a summary for existing file updates it."""
        store.register("collection", "doc.pdf", "Original summary")
        store.register("collection", "doc.pdf", "Updated summary")

        docs = store.get_all("collection")
        assert len(docs) == 1
        assert docs[0].summary == "Updated summary"

    def test_register_with_tags_roundtrip(self, store):
        """Test tags are JSON-persisted and decoded back into a list."""
        store.register("coll", "plan.pdf", "A floor plan.", tags=["Grundriss", "Brandschutz"])
        docs = store.get_all("coll")
        assert len(docs) == 1
        assert docs[0].summary == "A floor plan."
        assert docs[0].tags == ["Grundriss", "Brandschutz"]

    def test_register_without_tags_is_none(self, store):
        """Test a summary registered without tags decodes to None (not [])."""
        store.register("coll", "doc.pdf", "A summary.")
        docs = store.get_all("coll")
        assert docs[0].tags is None

    def test_register_tags_update_overwrites(self, store):
        """Test re-registering replaces the previously stored tags."""
        store.register("coll", "doc.pdf", "Summary.", tags=["Grundriss"])
        store.register("coll", "doc.pdf", "Summary.", tags=["Schnitt", "Schallschutz"])
        docs = store.get_all("coll")
        assert docs[0].tags == ["Schnitt", "Schallschutz"]

    @pytest.mark.asyncio
    async def test_get_all_async_with_tags(self, store):
        """Test async retrieval decodes the tags column."""
        store.register("acoll", "plan.pdf", "A plan.", tags=["Grundriss"])
        docs = await store.get_all_async("acoll")
        assert len(docs) == 1
        assert docs[0].tags == ["Grundriss"]

    def test_fresh_table_has_tags_column(self, temp_db):
        """A freshly-created document_metadata table includes the tags column."""
        from sqlalchemy import create_engine
        from sqlalchemy import inspect

        DocumentMetadataStore(temp_db)  # creates the table
        engine = create_engine(temp_db)
        columns = {c["name"] for c in inspect(engine).get_columns("document_metadata")}
        assert "tags" in columns
        engine.dispose()

    def test_existing_table_without_tags_is_migrated(self):
        """A pre-existing summaries table (no tags column) is renamed + migrated in place."""
        from sqlalchemy import create_engine
        from sqlalchemy import inspect
        from sqlalchemy import text

        with tempfile.TemporaryDirectory() as tmpdir:
            db_path = Path(tmpdir) / "legacy.db"
            db_url = f"sqlite:///{db_path}"

            # Simulate a legacy DB: create the table WITHOUT the tags column.
            engine = create_engine(db_url)
            with engine.connect() as conn:
                conn.execute(
                    text(
                        "CREATE TABLE summaries ("
                        "collection VARCHAR(256) NOT NULL, "
                        "filename VARCHAR(512) NOT NULL, "
                        "summary TEXT NOT NULL, "
                        "created_at DATETIME, "
                        "PRIMARY KEY (collection, filename))"
                    )
                )
                conn.execute(
                    text("INSERT INTO summaries (collection, filename, summary) VALUES ('c', 'old.pdf', 'Legacy.')")
                )
                conn.commit()

            assert "tags" not in {c["name"] for c in inspect(engine).get_columns("summaries")}

            # Constructing the store renames the legacy table and backfills columns.
            DocumentMetadataStore._tables_initialized.discard(db_url)
            store = DocumentMetadataStore(db_url)
            inspector = inspect(engine)
            assert not inspector.has_table("summaries")  # renamed away
            assert "tags" in {c["name"] for c in inspector.get_columns("document_metadata")}

            # Existing rows survive, tags default to None, new writes carry tags.
            docs = store.get_all("c")
            assert len(docs) == 1 and docs[0].tags is None
            store.register("c", "new.pdf", "New.", tags=["Grundriss"])
            new = {d.file_name: d for d in store.get_all("c")}
            assert new["new.pdf"].tags == ["Grundriss"]

            engine.dispose()
            DocumentMetadataStore.dispose_engine(db_url)

    def test_set_and_get_doc_class_roundtrip(self, store):
        """set_doc_class updates an existing row; get_doc_class reads it back."""
        store.register("coll", "doc.pdf", "A summary.")
        assert store.get_doc_class("coll", "doc.pdf") is None  # default null
        assert store.set_doc_class("coll", "doc.pdf", "oib_richtlinie") is True
        assert store.get_doc_class("coll", "doc.pdf") == "oib_richtlinie"
        # It also surfaces on the AvailableDocument read.
        docs = store.get_all("coll")
        assert docs[0].doc_class == "oib_richtlinie"

    def test_set_doc_class_without_row_returns_false(self, store):
        """set_doc_class never creates a row: no summary → False, still null."""
        assert store.set_doc_class("coll", "missing.pdf", "gesetz") is False
        assert store.get_doc_class("coll", "missing.pdf") is None

    def test_set_doc_class_clear(self, store):
        """A None value clears the stored doc_class back to null."""
        store.register("coll", "doc.pdf", "A summary.")
        store.set_doc_class("coll", "doc.pdf", "gesetz")
        assert store.set_doc_class("coll", "doc.pdf", None) is True
        assert store.get_doc_class("coll", "doc.pdf") is None

    def test_fresh_table_has_doc_class_column(self, temp_db):
        """A freshly-created document_metadata table includes the doc_class column."""
        from sqlalchemy import create_engine
        from sqlalchemy import inspect

        DocumentMetadataStore(temp_db)  # creates the table
        engine = create_engine(temp_db)
        columns = {c["name"] for c in inspect(engine).get_columns("document_metadata")}
        assert "doc_class" in columns
        engine.dispose()

    def test_existing_table_without_doc_class_is_migrated(self):
        """A pre-existing summaries table (no doc_class column) is migrated in place."""
        from sqlalchemy import create_engine
        from sqlalchemy import inspect
        from sqlalchemy import text

        with tempfile.TemporaryDirectory() as tmpdir:
            db_path = Path(tmpdir) / "legacy_dc.db"
            db_url = f"sqlite:///{db_path}"

            # Legacy DB: table with the tags column but WITHOUT doc_class.
            engine = create_engine(db_url)
            with engine.connect() as conn:
                conn.execute(
                    text(
                        "CREATE TABLE summaries ("
                        "collection VARCHAR(256) NOT NULL, "
                        "filename VARCHAR(512) NOT NULL, "
                        "summary TEXT NOT NULL, "
                        "tags TEXT, "
                        "created_at DATETIME, "
                        "PRIMARY KEY (collection, filename))"
                    )
                )
                conn.execute(
                    text("INSERT INTO summaries (collection, filename, summary) VALUES ('c', 'old.pdf', 'Legacy.')")
                )
                conn.commit()

            assert "doc_class" not in {c["name"] for c in inspect(engine).get_columns("summaries")}

            DocumentMetadataStore._tables_initialized.discard(db_url)
            store = DocumentMetadataStore(db_url)
            assert "doc_class" in {c["name"] for c in inspect(engine).get_columns("document_metadata")}

            # Existing rows survive with a null doc_class; new writes round-trip.
            assert store.get_doc_class("c", "old.pdf") is None
            store.set_doc_class("c", "old.pdf", "oib_leitfaden")
            assert store.get_doc_class("c", "old.pdf") == "oib_leitfaden"

            engine.dispose()
            DocumentMetadataStore.dispose_engine(db_url)

    def test_set_and_get_display_title_roundtrip(self, store):
        """set_display_title updates an existing row; get_display_title reads it back."""
        store.register("coll", "doc.pdf", "A summary.")
        assert store.get_display_title("coll", "doc.pdf") is None  # default null
        assert store.set_display_title("coll", "doc.pdf", "OIB-Richtlinie 2, Ausgabe Mai 2023") is True
        assert store.get_display_title("coll", "doc.pdf") == "OIB-Richtlinie 2, Ausgabe Mai 2023"
        # It also surfaces on the AvailableDocument read.
        docs = store.get_all("coll")
        assert docs[0].display_title == "OIB-Richtlinie 2, Ausgabe Mai 2023"

    def test_set_display_title_without_row_returns_false(self, store):
        """set_display_title never creates a row: no summary → False, still null."""
        assert store.set_display_title("coll", "missing.pdf", "Some title") is False
        assert store.get_display_title("coll", "missing.pdf") is None

    def test_set_display_title_clear(self, store):
        """A None value clears the override so the derived default applies again."""
        store.register("coll", "doc.pdf", "A summary.")
        store.set_display_title("coll", "doc.pdf", "A title")
        assert store.set_display_title("coll", "doc.pdf", None) is True
        assert store.get_display_title("coll", "doc.pdf") is None

    def test_display_titles_batch(self, store):
        """get_display_titles_batch returns only rows with a truthy stored title."""
        store.register("coll", "a.pdf", "A.")
        store.register("coll", "b.pdf", "B.")
        store.register("coll", "c.pdf", "C.")
        store.set_display_title("coll", "a.pdf", "Title A")
        store.set_display_title("coll", "c.pdf", "Title C")
        got = store.get_display_titles_batch("coll", ["a.pdf", "b.pdf", "c.pdf"])
        assert got == {"a.pdf": "Title A", "c.pdf": "Title C"}

    def test_fresh_table_has_display_title_column(self, temp_db):
        """A freshly-created document_metadata table includes the display_title column."""
        from sqlalchemy import create_engine
        from sqlalchemy import inspect

        DocumentMetadataStore(temp_db)  # creates the table
        engine = create_engine(temp_db)
        columns = {c["name"] for c in inspect(engine).get_columns("document_metadata")}
        assert "display_title" in columns
        engine.dispose()

    def test_legacy_summaries_table_renamed_preserving_rows(self):
        """A legacy ``summaries`` table is renamed to ``document_metadata`` in place,
        preserving existing rows and backfilling every added column."""
        from sqlalchemy import create_engine
        from sqlalchemy import inspect
        from sqlalchemy import text

        with tempfile.TemporaryDirectory() as tmpdir:
            db_url = f"sqlite:///{Path(tmpdir) / 'legacy_full.db'}"

            # Simulate the original schema: named ``summaries``, with the old index,
            # no doc_class/display_title columns.
            engine = create_engine(db_url)
            with engine.connect() as conn:
                conn.execute(
                    text(
                        "CREATE TABLE summaries ("
                        "collection VARCHAR(256) NOT NULL, "
                        "filename VARCHAR(512) NOT NULL, "
                        "summary TEXT NOT NULL, "
                        "tags TEXT, "
                        "created_at DATETIME, "
                        "PRIMARY KEY (collection, filename))"
                    )
                )
                conn.execute(text("CREATE INDEX idx_summaries_collection ON summaries (collection)"))
                conn.execute(
                    text(
                        "INSERT INTO summaries (collection, filename, summary) "
                        "VALUES ('oib_knowledge', 'oib-rl_2_ausgabe_mai_2023.pdf', 'Brandschutz.')"
                    )
                )
                conn.commit()

            DocumentMetadataStore._tables_initialized.discard(db_url)
            store = DocumentMetadataStore(db_url)

            inspector = inspect(engine)
            assert not inspector.has_table("summaries")  # legacy table renamed away
            assert inspector.has_table("document_metadata")
            columns = {c["name"] for c in inspector.get_columns("document_metadata")}
            assert {"tags", "doc_class", "display_title"} <= columns
            index_names = {ix["name"] for ix in inspector.get_indexes("document_metadata")}
            assert "idx_document_metadata_collection" in index_names
            assert "idx_summaries_collection" not in index_names  # legacy index dropped

            # The pre-existing row survived untouched and its new columns are null.
            docs = store.get_all("oib_knowledge")
            assert len(docs) == 1
            assert docs[0].file_name == "oib-rl_2_ausgabe_mai_2023.pdf"
            assert docs[0].summary == "Brandschutz."
            assert docs[0].doc_class is None
            assert docs[0].display_title is None

            engine.dispose()
            DocumentMetadataStore.dispose_engine(db_url)

    def test_get_all_empty_collection(self, store):
        """Test getting documents from empty collection returns empty list."""
        docs = store.get_all("nonexistent_collection")
        assert docs == []

    def test_get_all_different_collections(self, store):
        """Test documents are isolated by collection."""
        store.register("collection_a", "doc1.pdf", "Summary A1")
        store.register("collection_a", "doc2.pdf", "Summary A2")
        store.register("collection_b", "doc3.pdf", "Summary B1")

        docs_a = store.get_all("collection_a")
        docs_b = store.get_all("collection_b")

        assert len(docs_a) == 2
        assert len(docs_b) == 1
        assert docs_b[0].file_name == "doc3.pdf"

    def test_unregister_summary(self, store):
        """Test unregistering a document summary."""
        store.register("collection", "doc1.pdf", "Summary 1")
        store.register("collection", "doc2.pdf", "Summary 2")

        store.unregister("collection", "doc1.pdf")

        docs = store.get_all("collection")
        assert len(docs) == 1
        assert docs[0].file_name == "doc2.pdf"

    def test_unregister_nonexistent(self, store):
        """Test unregistering nonexistent document doesn't raise error."""
        store.unregister("collection", "nonexistent.pdf")  # Should not raise

    def test_clear_collection(self, store):
        """Test clearing all summaries in a collection."""
        store.register("collection", "doc1.pdf", "Summary 1")
        store.register("collection", "doc2.pdf", "Summary 2")
        store.register("other_collection", "doc3.pdf", "Summary 3")

        store.clear_collection("collection")

        assert store.get_all("collection") == []
        assert len(store.get_all("other_collection")) == 1

    def test_clear_all(self, store):
        """Test clearing all summaries across all collections."""
        store.register("collection1", "doc1.pdf", "Summary 1")
        store.register("collection2", "doc2.pdf", "Summary 2")

        store.clear_all()

        assert store.get_all("collection1") == []
        assert store.get_all("collection2") == []

    @pytest.mark.asyncio
    async def test_get_all_async(self, store):
        """Test async retrieval of document summaries."""
        store.register("async_collection", "doc1.pdf", "Async summary 1")
        store.register("async_collection", "doc2.pdf", "Async summary 2")

        docs = await store.get_all_async("async_collection")

        assert len(docs) == 2
        filenames = {doc.file_name for doc in docs}
        assert filenames == {"doc1.pdf", "doc2.pdf"}

    @pytest.mark.asyncio
    async def test_get_all_async_empty(self, store):
        """Test async retrieval from empty collection."""
        docs = await store.get_all_async("empty_collection")
        assert docs == []

    # -- update_tags -------------------------------------------------------

    def test_update_tags_roundtrip(self, store):
        """update_tags replaces tags on an existing row without touching summary."""
        store.register("coll", "plan.pdf", "A floor plan.", tags=["Grundriss"])

        assert store.update_tags("coll", "plan.pdf", ["Schnitt", "Brandschutz"]) is True

        docs = {d.file_name: d for d in store.get_all("coll")}
        assert docs["plan.pdf"].tags == ["Schnitt", "Brandschutz"]
        # Summary is never modified by a tag update.
        assert docs["plan.pdf"].summary == "A floor plan."

    def test_update_tags_on_row_without_tags(self, store):
        """A summary with NULL tags (legacy row) can be tagged via update_tags."""
        store.register("coll", "legacy.pdf", "Legacy summary.")

        assert store.update_tags("coll", "legacy.pdf", ["Bescheid"]) is True

        docs = {d.file_name: d for d in store.get_all("coll")}
        assert docs["legacy.pdf"].tags == ["Bescheid"]

    def test_update_tags_empty_clears(self, store):
        """An empty list clears the tags back to None (SQL NULL)."""
        store.register("coll", "plan.pdf", "A floor plan.", tags=["Grundriss"])

        assert store.update_tags("coll", "plan.pdf", []) is True

        docs = {d.file_name: d for d in store.get_all("coll")}
        assert docs["plan.pdf"].tags is None

    def test_update_tags_missing_row_returns_false(self, store):
        """Updating tags for a non-existent summary row returns False (no insert)."""
        assert store.update_tags("coll", "ghost.pdf", ["Grundriss"]) is False
        # No summary-less row is created (summary is NOT NULL).
        assert store.get_all("coll") == []

    # -- list_collections --------------------------------------------------

    def test_list_collections(self, store):
        """list_collections returns the distinct collections, sorted."""
        store.register("coll_b", "doc1.pdf", "S1")
        store.register("coll_a", "doc2.pdf", "S2")
        store.register("coll_a", "doc3.pdf", "S3")

        assert store.list_collections() == ["coll_a", "coll_b"]

    def test_list_collections_empty(self, store):
        """list_collections on an empty store returns an empty list."""
        assert store.list_collections() == []


# =============================================================================
# Factory Function Tests
# =============================================================================


class TestFactoryFunctions:
    """Tests for the factory module's summary registry functions."""

    @pytest.fixture(autouse=True)
    def reset_summary_store(self):
        """Reset the global summary store before each test."""
        from aiq_agent.knowledge import factory

        factory._document_metadata_store = None
        yield
        factory._document_metadata_store = None

    @pytest.fixture
    def temp_db_url(self):
        """Create a temporary SQLite database URL."""
        with tempfile.TemporaryDirectory() as tmpdir:
            db_path = Path(tmpdir) / "factory_test.db"
            yield f"sqlite:///{db_path}"
            DocumentMetadataStore.dispose_engine(f"sqlite:///{db_path}")

    def test_configure_summary_db(self, temp_db_url):
        """Test configuring the summary database."""
        from aiq_agent.knowledge import configure_summary_db

        configure_summary_db(temp_db_url)

        from aiq_agent.knowledge import factory

        assert factory._document_metadata_store is not None
        assert factory._document_metadata_store.db_url == temp_db_url

    def test_register_summary(self, temp_db_url):
        """Test registering a summary via factory function."""
        from aiq_agent.knowledge import configure_summary_db
        from aiq_agent.knowledge import get_available_documents
        from aiq_agent.knowledge.factory import register_summary

        configure_summary_db(temp_db_url)
        register_summary("test_collection", "test.pdf", "Test summary")

        docs = get_available_documents("test_collection")
        assert len(docs) == 1
        assert docs[0].file_name == "test.pdf"
        assert docs[0].summary == "Test summary"

    def test_register_summary_none_skipped(self, temp_db_url):
        """Test that None summaries are not registered."""
        from aiq_agent.knowledge import configure_summary_db
        from aiq_agent.knowledge import get_available_documents
        from aiq_agent.knowledge.factory import register_summary

        configure_summary_db(temp_db_url)
        register_summary("collection", "doc.pdf", None)

        docs = get_available_documents("collection")
        assert len(docs) == 0

    def test_register_summary_empty_string_skipped(self, temp_db_url):
        """Test that empty string summaries are not registered."""
        from aiq_agent.knowledge import configure_summary_db
        from aiq_agent.knowledge import get_available_documents
        from aiq_agent.knowledge.factory import register_summary

        configure_summary_db(temp_db_url)
        register_summary("collection", "doc.pdf", "")

        docs = get_available_documents("collection")
        assert len(docs) == 0

    def test_get_available_documents(self, temp_db_url):
        """Test getting available documents via factory function."""
        from aiq_agent.knowledge import configure_summary_db
        from aiq_agent.knowledge import get_available_documents
        from aiq_agent.knowledge.factory import register_summary

        configure_summary_db(temp_db_url)
        register_summary("my_collection", "file1.pdf", "Summary for file 1")
        register_summary("my_collection", "file2.pdf", "Summary for file 2")

        docs = get_available_documents("my_collection")
        assert len(docs) == 2
        assert all(isinstance(doc, AvailableDocument) for doc in docs)

    @pytest.mark.asyncio
    async def test_get_available_documents_async(self, temp_db_url):
        """Test async retrieval of available documents."""
        from aiq_agent.knowledge import configure_summary_db
        from aiq_agent.knowledge.factory import get_available_documents_async
        from aiq_agent.knowledge.factory import register_summary

        configure_summary_db(temp_db_url)
        register_summary("async_test", "doc.pdf", "Async test summary")

        docs = await get_available_documents_async("async_test")
        assert len(docs) == 1
        assert docs[0].file_name == "doc.pdf"

    def test_unregister_summary(self, temp_db_url):
        """Test unregistering a summary via factory function."""
        from aiq_agent.knowledge import configure_summary_db
        from aiq_agent.knowledge import get_available_documents
        from aiq_agent.knowledge.factory import register_summary
        from aiq_agent.knowledge.factory import unregister_summary

        configure_summary_db(temp_db_url)
        register_summary("collection", "doc1.pdf", "Summary 1")
        register_summary("collection", "doc2.pdf", "Summary 2")

        unregister_summary("collection", "doc1.pdf")

        docs = get_available_documents("collection")
        assert len(docs) == 1
        assert docs[0].file_name == "doc2.pdf"

    def test_clear_collection_summaries(self, temp_db_url):
        """Test clearing collection summaries via factory function."""
        from aiq_agent.knowledge import configure_summary_db
        from aiq_agent.knowledge import get_available_documents
        from aiq_agent.knowledge.factory import clear_collection_summaries
        from aiq_agent.knowledge.factory import register_summary

        configure_summary_db(temp_db_url)
        register_summary("coll1", "doc1.pdf", "Summary 1")
        register_summary("coll2", "doc2.pdf", "Summary 2")

        clear_collection_summaries("coll1")

        assert get_available_documents("coll1") == []
        assert len(get_available_documents("coll2")) == 1

    def test_clear_all_summaries(self, temp_db_url):
        """Test clearing all summaries via factory function."""
        from aiq_agent.knowledge import configure_summary_db
        from aiq_agent.knowledge import get_available_documents
        from aiq_agent.knowledge.factory import clear_all_summaries
        from aiq_agent.knowledge.factory import register_summary

        configure_summary_db(temp_db_url)
        register_summary("coll1", "doc1.pdf", "Summary 1")
        register_summary("coll2", "doc2.pdf", "Summary 2")

        clear_all_summaries()

        assert get_available_documents("coll1") == []
        assert get_available_documents("coll2") == []

    def test_lazy_initialization(self):
        """Test that summary store is lazily initialized with default DB."""
        from aiq_agent.knowledge import get_available_documents

        # Don't call configure_summary_db - should auto-initialize
        docs = get_available_documents("test_collection")
        assert docs == []  # Empty but works

        from aiq_agent.knowledge import factory

        assert factory._document_metadata_store is not None


# =============================================================================
# Reconciliation Backfill Tests
# =============================================================================


class TestReconcileCollectionSummaries:
    """Tests for the vector-index-vs-summaries-table reconciliation backstop.

    ``reconcile_collection_summaries`` only calls ``ingestor.list_files(...)``
    and, if present, an optional ``ingestor.get_document_text_sample(...)`` —
    so a duck-typed fake (not a real ``BaseIngestor`` subclass) is enough to
    exercise it in isolation, without any real vector store.
    """

    @pytest.fixture(autouse=True)
    def reset_summary_store(self):
        """Reset the global summary store before each test."""
        from aiq_agent.knowledge import factory

        factory._document_metadata_store = None
        yield
        factory._document_metadata_store = None

    @pytest.fixture
    def temp_db_url(self):
        """Create a temporary SQLite database URL."""
        with tempfile.TemporaryDirectory() as tmpdir:
            db_path = Path(tmpdir) / "reconcile_test.db"
            yield f"sqlite:///{db_path}"
            DocumentMetadataStore.dispose_engine(f"sqlite:///{db_path}")

    @staticmethod
    def _file_info(file_name, status=None):
        from aiq_agent.knowledge.schema import FileInfo
        from aiq_agent.knowledge.schema import FileStatus

        return FileInfo(
            file_id=file_name,
            file_name=file_name,
            collection_name="coll",
            status=status or FileStatus.SUCCESS,
        )

    class _FakeIngestor:
        """Stand-in for a backend ingestor with an indexed-file list and an
        optional text sampler."""

        def __init__(self, files, text_samples=None):
            self._files = files
            self._text_samples = text_samples or {}

        def list_files(self, collection_name):
            return self._files

        def get_document_text_sample(self, collection_name, file_name):
            return self._text_samples.get(file_name)

    def test_backfills_missing_document_with_text_sample(self, temp_db_url):
        """A document present in the index but absent from summaries gets a
        deterministic fallback summary derived from its sampled text."""
        from aiq_agent.knowledge import configure_summary_db
        from aiq_agent.knowledge import get_available_documents
        from aiq_agent.knowledge.factory import reconcile_collection_summaries

        configure_summary_db(temp_db_url)

        ingestor = self._FakeIngestor(
            files=[self._file_info("orphan.pdf")],
            text_samples={"orphan.pdf": "Dies ist ein Brandschutzkonzept fuer das Gebaeude."},
        )

        backfilled = reconcile_collection_summaries(ingestor, "coll")

        assert backfilled == 1
        docs = {d.file_name: d for d in get_available_documents("coll")}
        assert "orphan.pdf" in docs
        assert "Brandschutzkonzept" in docs["orphan.pdf"].summary

    def test_no_sampler_falls_back_to_filename(self, temp_db_url):
        """Backends without the optional text-sampling hook still get a row —
        visibility beats a perfect summary."""
        from aiq_agent.knowledge import configure_summary_db
        from aiq_agent.knowledge import get_available_documents
        from aiq_agent.knowledge.factory import reconcile_collection_summaries

        configure_summary_db(temp_db_url)

        class NoSamplerIngestor:
            def list_files(self, collection_name):
                return [TestReconcileCollectionSummaries._file_info("mystery.pdf")]

        backfilled = reconcile_collection_summaries(NoSamplerIngestor(), "coll")

        assert backfilled == 1
        docs = {d.file_name: d for d in get_available_documents("coll")}
        assert docs["mystery.pdf"].summary == "Indexed document: mystery.pdf"

    def test_noop_when_already_consistent(self, temp_db_url):
        """A document already present in the summaries table is left alone."""
        from aiq_agent.knowledge import configure_summary_db
        from aiq_agent.knowledge import get_available_documents
        from aiq_agent.knowledge.factory import reconcile_collection_summaries
        from aiq_agent.knowledge.factory import register_summary

        configure_summary_db(temp_db_url)
        register_summary("coll", "already.pdf", "Already has a summary.")

        ingestor = self._FakeIngestor(files=[self._file_info("already.pdf")])

        backfilled = reconcile_collection_summaries(ingestor, "coll")

        assert backfilled == 0
        docs = get_available_documents("coll")
        assert len(docs) == 1
        assert docs[0].summary == "Already has a summary."

    def test_ignores_files_not_indexed_successfully(self, temp_db_url):
        """A FAILED file was never actually indexed, so it must never get a
        fallback summary row."""
        from aiq_agent.knowledge import configure_summary_db
        from aiq_agent.knowledge import get_available_documents
        from aiq_agent.knowledge.factory import reconcile_collection_summaries
        from aiq_agent.knowledge.schema import FileStatus

        configure_summary_db(temp_db_url)

        ingestor = self._FakeIngestor(files=[self._file_info("broken.pdf", status=FileStatus.FAILED)])

        backfilled = reconcile_collection_summaries(ingestor, "coll")

        assert backfilled == 0
        assert get_available_documents("coll") == []

    def test_sampler_exception_falls_back_to_filename(self, temp_db_url):
        """A broken ``get_document_text_sample`` must not break reconciliation
        (fail-open, same contract as the primary summary path)."""
        from aiq_agent.knowledge import configure_summary_db
        from aiq_agent.knowledge import get_available_documents
        from aiq_agent.knowledge.factory import reconcile_collection_summaries

        configure_summary_db(temp_db_url)

        class BrokenSamplerIngestor:
            def list_files(self, collection_name):
                return [TestReconcileCollectionSummaries._file_info("weird.pdf")]

            def get_document_text_sample(self, collection_name, file_name):
                raise RuntimeError("chroma unreachable")

        backfilled = reconcile_collection_summaries(BrokenSamplerIngestor(), "coll")

        assert backfilled == 1
        docs = {d.file_name: d for d in get_available_documents("coll")}
        assert docs["weird.pdf"].summary == "Indexed document: weird.pdf"

    def test_list_files_exception_returns_zero(self, temp_db_url):
        """A backend that cannot list files at all fails open (never raises)."""
        from aiq_agent.knowledge import configure_summary_db
        from aiq_agent.knowledge.factory import reconcile_collection_summaries

        configure_summary_db(temp_db_url)

        class BrokenListFilesIngestor:
            def list_files(self, collection_name):
                raise RuntimeError("chroma down")

        assert reconcile_collection_summaries(BrokenListFilesIngestor(), "coll") == 0


# =============================================================================
# Integration Tests
# =============================================================================


class TestSummaryIntegration:
    """Integration tests for the summary storage workflow."""

    @pytest.fixture
    def temp_db_url(self):
        """Create a temporary SQLite database URL."""
        with tempfile.TemporaryDirectory() as tmpdir:
            db_path = Path(tmpdir) / "integration_test.db"
            yield f"sqlite:///{db_path}"
            DocumentMetadataStore.dispose_engine(f"sqlite:///{db_path}")

    @pytest.fixture(autouse=True)
    def reset_store(self):
        """Reset global store before each test."""
        from aiq_agent.knowledge import factory

        factory._document_metadata_store = None
        yield
        factory._document_metadata_store = None

    def test_full_workflow(self, temp_db_url):
        """Test complete summary storage workflow."""
        from aiq_agent.knowledge import configure_summary_db
        from aiq_agent.knowledge import get_available_documents
        from aiq_agent.knowledge.factory import clear_collection_summaries
        from aiq_agent.knowledge.factory import register_summary
        from aiq_agent.knowledge.factory import unregister_summary

        # 1. Configure database
        configure_summary_db(temp_db_url)

        # 2. Simulate document ingestion with summaries
        collection = "session_12345"
        register_summary(collection, "report.pdf", "Annual financial report for 2024.")
        register_summary(collection, "data.csv", "Sales data spreadsheet.")
        register_summary(collection, "notes.txt", "Meeting notes from Q4 review.")

        # 3. Retrieve for agent context
        docs = get_available_documents(collection)
        assert len(docs) == 3

        # 4. User deletes a file
        unregister_summary(collection, "notes.txt")
        docs = get_available_documents(collection)
        assert len(docs) == 2

        # 5. Session ends, cleanup
        clear_collection_summaries(collection)
        docs = get_available_documents(collection)
        assert len(docs) == 0

    def test_multiple_collections_isolation(self, temp_db_url):
        """Test that different sessions/collections are isolated."""
        from aiq_agent.knowledge import configure_summary_db
        from aiq_agent.knowledge import get_available_documents
        from aiq_agent.knowledge.factory import register_summary

        configure_summary_db(temp_db_url)

        # Two different user sessions
        session_a = "user_alice_session_001"
        session_b = "user_bob_session_002"

        register_summary(session_a, "alice_doc.pdf", "Alice's document")
        register_summary(session_b, "bob_doc.pdf", "Bob's document")
        register_summary(session_b, "bob_doc2.pdf", "Bob's second document")

        # Each user only sees their own documents
        alice_docs = get_available_documents(session_a)
        bob_docs = get_available_documents(session_b)

        assert len(alice_docs) == 1
        assert alice_docs[0].file_name == "alice_doc.pdf"

        assert len(bob_docs) == 2
        assert {d.file_name for d in bob_docs} == {"bob_doc.pdf", "bob_doc2.pdf"}

    @pytest.mark.asyncio
    async def test_mixed_sync_async_operations(self, temp_db_url):
        """Test mixing sync registration with async retrieval."""
        from aiq_agent.knowledge import configure_summary_db
        from aiq_agent.knowledge.factory import get_available_documents_async
        from aiq_agent.knowledge.factory import register_summary

        configure_summary_db(temp_db_url)

        # Sync registration (as backend adapters do)
        register_summary("mixed_test", "doc1.pdf", "Document 1 summary")
        register_summary("mixed_test", "doc2.pdf", "Document 2 summary")

        # Async retrieval (as agent code does)
        docs = await get_available_documents_async("mixed_test")

        assert len(docs) == 2
        assert all(doc.summary is not None for doc in docs)


# =============================================================================
# Migration-failure caching guard
# =============================================================================


class TestMigrationFailureNotCached:
    """A failed column migration must NOT mark the store initialized, so the next
    access retries it instead of writing against a missing column. The legacy
    ``summaries`` table is renamed before columns are added, so a first-pass
    failure leaves a renamed-but-incomplete table the retry completes."""

    @staticmethod
    def _legacy_db(dir_path):
        """Create a pre-columns ``summaries`` table and return (url, engine)."""
        from sqlalchemy import create_engine
        from sqlalchemy import text

        db_url = f"sqlite:///{dir_path / 'legacy.db'}"
        engine = create_engine(db_url)
        with engine.connect() as conn:
            conn.execute(
                text(
                    "CREATE TABLE summaries ("
                    "collection VARCHAR(256) NOT NULL, "
                    "filename VARCHAR(512) NOT NULL, "
                    "summary TEXT NOT NULL, "
                    "created_at DATETIME, "
                    "PRIMARY KEY (collection, filename))"
                )
            )
            conn.commit()
        return db_url, engine

    @staticmethod
    def _columns(engine):
        from sqlalchemy import inspect

        inspector = inspect(engine)
        table = "document_metadata" if inspector.has_table("document_metadata") else "summaries"
        return {c["name"] for c in inspector.get_columns(table)}

    @staticmethod
    def _flaky_add_column(calls):
        """Return an ``_add_column_if_missing`` that raises on its first call."""
        real = DocumentMetadataStore._add_column_if_missing

        def flaky(self, conn, column, ddl_type="TEXT"):
            calls["n"] += 1
            if calls["n"] == 1:
                raise RuntimeError("simulated ALTER failure")
            return real(self, conn, column, ddl_type)

        return flaky

    def test_failed_sync_migration_not_cached_then_retried(self, monkeypatch):
        with tempfile.TemporaryDirectory() as tmpdir:
            db_url, engine = self._legacy_db(Path(tmpdir))
            DocumentMetadataStore._tables_initialized.discard(db_url)
            try:
                calls = {"n": 0}
                monkeypatch.setattr(DocumentMetadataStore, "_add_column_if_missing", self._flaky_add_column(calls))

                # First construction: migration "fails" -> not cached, no tags col.
                DocumentMetadataStore(db_url)
                assert db_url not in DocumentMetadataStore._tables_initialized
                assert "tags" not in self._columns(engine)

                # Next construction retries the migration -> succeeds and caches.
                DocumentMetadataStore(db_url)
                assert db_url in DocumentMetadataStore._tables_initialized
                assert "tags" in self._columns(engine)
            finally:
                engine.dispose()
                DocumentMetadataStore.dispose_engine(db_url)
                DocumentMetadataStore._tables_initialized.discard(db_url)

    @pytest.mark.asyncio
    async def test_failed_async_migration_not_cached_then_retried(self, monkeypatch):
        with tempfile.TemporaryDirectory() as tmpdir:
            db_url, engine = self._legacy_db(Path(tmpdir))
            DocumentMetadataStore._tables_initialized.discard(db_url)
            try:
                calls = {"n": 0}
                monkeypatch.setattr(DocumentMetadataStore, "_add_column_if_missing", self._flaky_add_column(calls))

                # First ensure: migration "fails" -> not cached, no tags col.
                await DocumentMetadataStore._ensure_table_async(db_url)
                assert db_url not in DocumentMetadataStore._tables_initialized
                assert "tags" not in self._columns(engine)

                # Next ensure retries the migration -> succeeds and caches.
                await DocumentMetadataStore._ensure_table_async(db_url)
                assert db_url in DocumentMetadataStore._tables_initialized
                assert "tags" in self._columns(engine)
            finally:
                engine.dispose()
                await DocumentMetadataStore.dispose_engine_async(db_url)
                DocumentMetadataStore._tables_initialized.discard(db_url)


# =============================================================================
# Text Extraction Tests (Foundational RAG multi-format support)
# =============================================================================


class TestExtractText:
    """Tests for the _extract_text helper that supports PDF, DOCX, PPTX, TXT, and MD."""

    def test_txt_extraction(self):
        """Test text extraction from a .txt file."""
        from knowledge_layer.foundational_rag.adapter import _extract_text

        with tempfile.NamedTemporaryFile(suffix=".txt", mode="w", delete=False, encoding="utf-8") as f:
            f.write("Hello world. This is a test document for summarization.")
            f.flush()
            result = _extract_text(f.name)

        assert result is not None
        assert "Hello world" in result

    def test_md_extraction(self):
        """Test text extraction from a .md file."""
        from knowledge_layer.foundational_rag.adapter import _extract_text

        with tempfile.NamedTemporaryFile(suffix=".md", mode="w", delete=False, encoding="utf-8") as f:
            f.write("# Heading\n\nSome markdown content with **bold** text.")
            f.flush()
            result = _extract_text(f.name)

        assert result is not None
        assert "Heading" in result
        assert "bold" in result

    def test_txt_max_chars_truncation(self):
        """Test that text extraction respects max_chars limit."""
        from knowledge_layer.foundational_rag.adapter import _extract_text

        with tempfile.NamedTemporaryFile(suffix=".txt", mode="w", delete=False, encoding="utf-8") as f:
            f.write("A" * 10000)
            f.flush()
            result = _extract_text(f.name, max_chars=100)

        assert result is not None
        assert len(result) <= 100

    def test_empty_txt_returns_none(self):
        """Test that empty files return None."""
        from knowledge_layer.foundational_rag.adapter import _extract_text

        with tempfile.NamedTemporaryFile(suffix=".txt", mode="w", delete=False, encoding="utf-8") as f:
            f.write("")
            f.flush()
            result = _extract_text(f.name)

        assert result is None

    def test_unsupported_extension_returns_none(self):
        """Test that unsupported file extensions return None."""
        from knowledge_layer.foundational_rag.adapter import _extract_text

        with tempfile.NamedTemporaryFile(suffix=".csv", mode="w", delete=False, encoding="utf-8") as f:
            f.write("col1,col2\nval1,val2")
            f.flush()
            result = _extract_text(f.name)

        assert result is None

    def test_nonexistent_file_returns_none(self):
        """Test that a nonexistent file returns None gracefully."""
        from knowledge_layer.foundational_rag.adapter import _extract_text

        result = _extract_text("/nonexistent/path/file.txt")
        assert result is None

    def test_pdf_extraction(self):
        """Test PDF text extraction using pypdf."""
        from knowledge_layer.foundational_rag.adapter import _extract_text

        pytest.importorskip("pypdf")
        pdf_dir = Path(__file__).parent / "data"
        pdf_files = list(pdf_dir.glob("*.pdf"))
        if not pdf_files:
            pytest.skip("No test PDFs in data/ directory")

        result = _extract_text(str(pdf_files[0]))
        assert result is not None
        assert len(result) > 0

    def test_docx_extraction(self):
        """Test DOCX text extraction using docx2txt."""
        from knowledge_layer.foundational_rag.adapter import _extract_text

        pytest.importorskip("docx2txt")

        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
            tmp_path = f.name

        try:
            from docx import Document

            doc = Document()
            doc.add_paragraph("Test paragraph for DOCX extraction.")
            doc.save(tmp_path)

            result = _extract_text(tmp_path)
            assert result is not None
            assert "Test paragraph" in result
        except ImportError:
            pytest.skip("python-docx not installed for DOCX creation")

    def test_pptx_extraction(self):
        """Test PPTX text extraction using python-pptx."""
        from knowledge_layer.foundational_rag.adapter import _extract_text

        pytest.importorskip("pptx")
        from pptx import Presentation

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            tmp_path = f.name

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Test Slide Title"
        slide.placeholders[1].text = "Slide body content for testing."
        prs.save(tmp_path)

        result = _extract_text(tmp_path)
        assert result is not None
        assert "Test Slide Title" in result


# =============================================================================
# File Summary Generation Tests
# =============================================================================


class TestGenerateFileSummary:
    """Tests for _generate_file_summary with multi-format support."""

    def test_no_llm_returns_none(self):
        """Test that None LLM returns None."""
        from knowledge_layer.foundational_rag.adapter import _generate_file_summary

        with tempfile.NamedTemporaryFile(suffix=".txt", mode="w", delete=False, encoding="utf-8") as f:
            f.write("Some content")
            f.flush()
            result = _generate_file_summary(f.name, llm=None)

        assert result is None

    def test_unsupported_format_returns_none(self):
        """Test that unsupported file formats return None even with LLM."""
        from unittest.mock import MagicMock

        from knowledge_layer.foundational_rag.adapter import _generate_file_summary

        mock_llm = MagicMock()
        with tempfile.NamedTemporaryFile(suffix=".csv", mode="w", delete=False, encoding="utf-8") as f:
            f.write("data")
            f.flush()
            result = _generate_file_summary(f.name, llm=mock_llm)

        assert result is None
        mock_llm.invoke.assert_not_called()

    def test_txt_summary_calls_llm(self):
        """Test that a .txt file triggers LLM summarization."""
        from unittest.mock import MagicMock

        from knowledge_layer.foundational_rag.adapter import _generate_file_summary

        mock_llm = MagicMock()
        mock_llm.invoke.return_value = MagicMock(content="A one-sentence summary.")

        with tempfile.NamedTemporaryFile(suffix=".txt", mode="w", delete=False, encoding="utf-8") as f:
            f.write("This is a long document about artificial intelligence and machine learning.")
            f.flush()
            result = _generate_file_summary(f.name, llm=mock_llm)

        assert result == "A one-sentence summary."
        mock_llm.invoke.assert_called_once()

    def test_md_summary_calls_llm(self):
        """Test that a .md file triggers LLM summarization."""
        from unittest.mock import MagicMock

        from knowledge_layer.foundational_rag.adapter import _generate_file_summary

        mock_llm = MagicMock()
        mock_llm.invoke.return_value = MagicMock(content="Markdown summary.")

        with tempfile.NamedTemporaryFile(suffix=".md", mode="w", delete=False, encoding="utf-8") as f:
            f.write("# Research Notes\n\nKey findings about neural networks.")
            f.flush()
            result = _generate_file_summary(f.name, llm=mock_llm)

        assert result == "Markdown summary."
        mock_llm.invoke.assert_called_once()

    def test_llm_failure_returns_none(self):
        """Test graceful handling when LLM raises an exception."""
        from unittest.mock import MagicMock

        from knowledge_layer.foundational_rag.adapter import _generate_file_summary

        mock_llm = MagicMock()
        mock_llm.invoke.side_effect = RuntimeError("LLM API error")

        with tempfile.NamedTemporaryFile(suffix=".txt", mode="w", delete=False, encoding="utf-8") as f:
            f.write("Some document content.")
            f.flush()
            result = _generate_file_summary(f.name, llm=mock_llm)

        assert result is None

    def test_summarizable_extensions_constant(self):
        """Test that SUMMARIZABLE_EXTENSIONS contains expected formats."""
        from knowledge_layer.foundational_rag.adapter import SUMMARIZABLE_EXTENSIONS

        assert ".pdf" in SUMMARIZABLE_EXTENSIONS
        assert ".docx" in SUMMARIZABLE_EXTENSIONS
        assert ".pptx" in SUMMARIZABLE_EXTENSIONS
        assert ".txt" in SUMMARIZABLE_EXTENSIONS
        assert ".md" in SUMMARIZABLE_EXTENSIONS
        assert ".csv" not in SUMMARIZABLE_EXTENSIONS
